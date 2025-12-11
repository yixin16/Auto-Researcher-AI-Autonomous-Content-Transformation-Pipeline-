# ppt_generator.py

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO
import logging
from pathlib import Path
import matplotlib.pyplot as plt
import re
from typing import List, Dict, Optional, Tuple

logger = logging.getLogger(__name__)
SLIDE_DIR = Path("outputs/slides")
SLIDE_DIR.mkdir(parents=True, exist_ok=True)


class DesignTheme:
    """
    Aesthetic: 'Soft Modern Minimalist'
    Palette: Slate Blue, Sage Green, Warm Gray, Off-White.
    """
    
    # Palette
    PRIMARY = RGBColor(47, 53, 66)        # Soft Charcoal/Slate (Text & Headers)
    ACCENT_1 = RGBColor(96, 163, 188)     # Dusty Blue (Highlights/Questions)
    ACCENT_2 = RGBColor(120, 224, 143)    # Soft Sage (Success/Points)
    BG_COLOR = RGBColor(241, 242, 246)    # Very Light Warm Gray (Slide Background)
    CARD_BG = RGBColor(255, 255, 255)     # White (Content Cards)
    TEXT_MAIN = RGBColor(87, 96, 111)     # Soft Dark Gray (Body Text)
    
    # Fonts
    TITLE_FONT = "Segoe UI Light"  # Clean, thin, modern
    BODY_FONT = "Segoe UI"         # Readable
    
    # Sizes (Requested Constraints)
    FONT_SIZE_TITLE = 20           # For Section Headers
    FONT_SIZE_BODY = 15            # For Content
    
    # Layout
    MARGIN = Inches(0.5)
    HEADER_HEIGHT = Inches(1.0)
    FOOTER_HEIGHT = Inches(0.4)


class PPTGenerator:
    """Aesthetic PowerPoint generator with card-based layouts"""
    
    def __init__(self):
        self.prs = Presentation()
        # 16:9 Widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = DesignTheme()
        self.slide_counter = 0
    
    def _set_background(self, slide):
        """Sets a soft off-white background for the whole slide"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.BG_COLOR

    def _add_card_background(self, slide, left, top, width, height, color=None):
        """Adds a rounded card with a very subtle shadow"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color if color else self.theme.CARD_BG
        
        # Remove outline for cleaner look
        shape.line.fill.background()
        
        # Soft Shadow
        shadow = shape.shadow
        shadow.inherit = False
        shadow.blur_radius = Pt(8)
        shadow.distance = Pt(3)
        shadow.angle = 90
        shadow.transparency = 0.85 # Very subtle
        return shape

    def _add_footer(self, slide, text: str = None):
        """Minimalist footer"""
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, self.prs.slide_height - self.theme.FOOTER_HEIGHT,
            self.prs.slide_width, self.theme.FOOTER_HEIGHT
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = self.theme.PRIMARY
        footer.line.fill.background()
        
        if text:
            tf = footer.text_frame
            tf.text = text
            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.RIGHT
            p.font.name = self.theme.BODY_FONT

    def _add_header(self, slide, title: str):
        """Clean, minimal header"""
        # Title Text
        title_box = slide.shapes.add_textbox(
            self.theme.MARGIN, Inches(0.2),
            self.prs.slide_width - self.theme.MARGIN*2, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme.TITLE_FONT
        p.font.size = Pt(32) 
        p.font.color.rgb = self.theme.PRIMARY
        p.alignment = PP_ALIGN.LEFT
        
        # Subtle divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.theme.MARGIN, Inches(1.0),
            Inches(12.33), Inches(0.01) # Very thin line
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.ACCENT_1
        line.line.fill.background()

    def _style_text_box(self, text_frame, text, font_size=None, is_bold=False, color=None):
        """Applies the requested aesthetic styling"""
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = text
        p.font.name = self.theme.BODY_FONT
        p.font.size = Pt(font_size if font_size else self.theme.FONT_SIZE_BODY)
        p.font.color.rgb = color if color else self.theme.TEXT_MAIN
        p.font.bold = is_bold
        p.line_spacing = 1.2

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Elegant Title Slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        # Center Card
        card_w, card_h = Inches(10), Inches(5)
        left = (self.prs.slide_width - card_w) / 2
        top = (self.prs.slide_height - card_h) / 2
        self._add_card_background(slide, left, top, card_w, card_h)
        
        # Title
        title_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(1), card_w - Inches(1), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme.TITLE_FONT
        p.font.size = Pt(44)
        p.font.color.rgb = self.theme.PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(left + Inches(1), top + Inches(3), card_w - Inches(2), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.name = self.theme.BODY_FONT
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme.ACCENT_1
            p.alignment = PP_ALIGN.CENTER

        # Footer Date
        import datetime
        date_str = datetime.datetime.now().strftime('%B %d, %Y')
        date_box = slide.shapes.add_textbox(0, Inches(6.8), self.prs.slide_width, Inches(0.5))
        p = date_box.text_frame.paragraphs[0]
        p.text = date_str
        p.font.size = Pt(12)
        p.font.color.rgb = self.theme.TEXT_MAIN
        p.alignment = PP_ALIGN.CENTER

        self.slide_counter += 1

    def add_visual_slide(self, header: str, summary: str, image_url: Optional[str] = None):
        """Clean Visual Slide with Card Layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, header)
        
        # Card Background
        self._add_card_background(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5))
        
        # Text (Left)
        textbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(6.5), Inches(4.8)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        
        if len(summary) > 800: summary = summary[:800] + "..."
        self._style_text_box(tf, summary, font_size=self.theme.FONT_SIZE_BODY)
        
        # Image (Right)
        if image_url:
            try:
                resp = requests.get(image_url, timeout=5)
                img_stream = BytesIO(resp.content)
                pic = slide.shapes.add_picture(
                    img_stream, Inches(7.5), Inches(1.6), width=Inches(5.0), height=Inches(4.9)
                )
                pic.line.color.rgb = self.theme.ACCENT_1
                pic.line.width = Pt(1)
            except Exception:
                pass
        
        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def add_analysis_slide(self, header: str, key_points: List[str], insights: List[str]):
        """Two-Card Analysis Layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, header)
        
        # === Left Card: Key Facts ===
        self._add_card_background(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
        
        lt_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.5))
        self._style_text_box(lt_box.text_frame, "Key Facts", font_size=self.theme.FONT_SIZE_TITLE, is_bold=True, color=self.theme.ACCENT_1)
        
        lc_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.4), Inches(4.5))
        tf = lc_box.text_frame
        tf.word_wrap = True
        
        for point in key_points[:5]:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(self.theme.FONT_SIZE_BODY)
            p.font.name = self.theme.BODY_FONT
            p.font.color.rgb = self.theme.TEXT_MAIN
            p.space_after = Pt(10)

        # === Right Card: Insights ===
        self._add_card_background(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5))
        
        rt_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.5))
        self._style_text_box(rt_box.text_frame, "Strategic Insights", font_size=self.theme.FONT_SIZE_TITLE, is_bold=True, color=self.theme.ACCENT_2)
        
        rc_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.4), Inches(4.5))
        tf = rc_box.text_frame
        tf.word_wrap = True
        
        for insight in insights[:4]:
            p = tf.add_paragraph()
            p.text = f"→ {insight}"
            p.font.size = Pt(self.theme.FONT_SIZE_BODY)
            p.font.name = self.theme.BODY_FONT
            p.font.color.rgb = self.theme.TEXT_MAIN
            p.space_after = Pt(12)

        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def add_bullet_slide(self, title: str, items: List[str], icon: str = "•"):
        """Clean Single Card List"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, title)
        
        self._add_card_background(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5))
        
        textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.33), Inches(5.0))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        display_items = items[:7]
        for item in display_items:
            p = tf.add_paragraph()
            p.text = f"{icon}  {item}"
            p.font.size = Pt(16) if len(display_items) < 5 else Pt(15)
            p.font.name = self.theme.BODY_FONT
            p.font.color.rgb = self.theme.TEXT_MAIN
            p.space_after = Pt(14)
            
        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def create_chart_image(self, data: Dict, suffix: str) -> Optional[str]:
        """Aesthetic Matplotlib Chart"""
        try:
            plt.style.use('seaborn-v0_8-ticks')
            fig, ax = plt.subplots(figsize=(10, 6))
            
            colors = ['#60a3bc', '#78e08f', '#f6b93b', '#e55039', '#4a69bd']
            
            bars = ax.bar(
                data['labels'], 
                data['values'], 
                color=colors[:len(data['labels'])], 
                zorder=3,
                alpha=0.9
            )
            
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_visible(False)
            ax.spines['bottom'].set_color('#DDDDDD')
            
            ax.grid(axis='y', linestyle='-', alpha=0.3, color='#DDDDDD', zorder=0)
            
            ax.set_title(data.get('title', 'Data Analysis'), fontsize=16, color='#2f3542', pad=20)
            ax.tick_params(axis='x', colors='#57606f', labelsize=11)
            ax.tick_params(axis='y', colors='#57606f', labelsize=11)
            
            plt.tight_layout()
            path = SLIDE_DIR / f"chart_{suffix}.png"
            plt.savefig(path, dpi=150, bbox_inches='tight', transparent=False)
            plt.close()
            return str(path)
        except Exception as e:
            logger.error(f"Chart error: {e}")
            return None

    def add_chart_slide(self, header: str, chart_path: str):
        """Centered Chart in Card"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, header)
        
        self._add_card_background(slide, Inches(1.0), Inches(1.3), Inches(11.33), Inches(5.5))
        
        try:
            slide.shapes.add_picture(chart_path, Inches(2.2), Inches(1.5), width=Inches(9))
        except Exception:
            pass
            
        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def add_qna_slide(self, qna_list: List[Tuple[str, str]]):
        """
        Interactive Q&A slide with 'Soft Modern' aesthetic.
        Q: Dusty Blue Card (Accent 1)
        A: White Card
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, "Discussion & Q&A")
        
        y_pos = Inches(1.6)
        
        for i, (q, a) in enumerate(qna_list[:2]):
            
            # --- Question Card ---
            self._add_card_background(
                slide, Inches(1.0), y_pos, Inches(11.33), Inches(0.8), 
                color=self.theme.ACCENT_1
            )
            
            q_tb = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.15), Inches(10.9), Inches(0.5))
            q_p = q_tb.text_frame.paragraphs[0]
            q_p.text = f"Q: {q}"
            q_p.font.name = self.theme.TITLE_FONT
            q_p.font.size = Pt(16)
            q_p.font.bold = True
            q_p.font.color.rgb = RGBColor(255, 255, 255)
            
            y_pos += Inches(0.9)
            
            # --- Answer Card ---
            self._add_card_background(
                slide, Inches(1.5), y_pos, Inches(10.83), Inches(1.5)
            )
            
            a_tb = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.1), Inches(10.4), Inches(1.3))
            tf = a_tb.text_frame
            tf.word_wrap = True
            
            if len(a) > 400: a = a[:400] + "..."
            self._style_text_box(tf, f"A: {a}", font_size=self.theme.FONT_SIZE_BODY)
            
            y_pos += Inches(1.8)
            
        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def add_swot_slide(self, swot: Dict):
        """Clean 4-Quadrant SWOT"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_header(slide, "Strategic SWOT Analysis")
        
        margin_x = Inches(0.8)
        margin_y = Inches(1.5)
        box_w = Inches(5.7)
        box_h = Inches(2.6)
        gap = Inches(0.3)
        
        def make_quadrant(title, text, x, y, accent_color):
            self._add_card_background(slide, x, y, box_w, box_h)
            
            tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), box_w, Inches(0.5))
            self._style_text_box(tb.text_frame, title, font_size=self.theme.FONT_SIZE_TITLE, is_bold=True, color=accent_color)
            
            cb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), box_w - Inches(0.4), box_h - Inches(0.9))
            tf = cb.text_frame
            tf.word_wrap = True
            
            if len(text) > 300: text = text[:300] + "..."
            self._style_text_box(tf, text, font_size=self.theme.FONT_SIZE_BODY)

        make_quadrant("STRENGTHS", swot.get('S', ''), margin_x, margin_y, self.theme.ACCENT_2)
        make_quadrant("WEAKNESSES", swot.get('W', ''), margin_x + box_w + gap, margin_y, RGBColor(235, 77, 75))
        make_quadrant("OPPORTUNITIES", swot.get('O', ''), margin_x, margin_y + box_h + gap, self.theme.ACCENT_1)
        make_quadrant("THREATS", swot.get('T', ''), margin_x + box_w + gap, margin_y + box_h + gap, RGBColor(240, 147, 43))

        self._add_footer(slide, f"{self.slide_counter}")
        self.slide_counter += 1

    def save(self, filename: str) -> Path:
        clean_name = re.sub(r'[\\/*?:"<>|]', "", filename)[:50]
        path = SLIDE_DIR / f"{clean_name}.pptx"
        self.prs.save(path)
        logger.info(f"💾 Presentation saved: {path}")
        return path